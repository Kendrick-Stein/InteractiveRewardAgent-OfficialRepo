from __future__ import annotations

import json
import os
import re
import shlex
import subprocess
import sys
import tempfile
import time
from dataclasses import dataclass
from typing import Any, Dict, Optional, Tuple

from smolagents import Tool

API_URL = (os.getenv("IRA_BASE_URL") or os.getenv("OPENAI_BASE_URL") or "https://api.cometapi.com/v1/").rstrip("/") + "/"
MODEL_NAME = "o3"

GENERATED_CHECKER_PATH = os.path.join(tempfile.gettempdir(), f"ira_ppt_checker_{os.getpid()}.py")
HTTP_RETRY = 3
HTTP_TIMEOUT_SEC = 120

# ==== Auto-injected extraction helpers for PPTX (copied from ppt-test/exame.py, trimmed) ====
EXTRACT_PRELUDE = r"""
# ==== Auto-injected extraction helpers (from extractpptinfo.py, simplified) ====
from typing import Any, Dict, List, Optional

try:
    from PIL import Image  # Optional: for reading original image pixel sizes
    PIL_AVAILABLE = True
except Exception:
    PIL_AVAILABLE = False

from pptx import Presentation  # safe to import here; generated script may also import
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU_PER_INCH = 914400
POINTS_PER_INCH = 72

def emu_to_in(emu: Optional[int]) -> Optional[float]:
    if emu is None:
        return None
    try:
        return round(emu / EMU_PER_INCH, 4)
    except Exception:
        return None

def emu_to_pt(emu: Optional[int]) -> Optional[float]:
    if emu is None:
        return None
    try:
        return round(emu / EMU_PER_INCH * POINTS_PER_INCH, 2)
    except Exception:
        return None

def rgb_color_to_hex(color) -> Optional[str]:
    try:
        if color is None:
            return None
        rgb = getattr(color, "rgb", None)
        if rgb is not None:
            return str(rgb)
        theme_color = getattr(color, "theme_color", None)
        if theme_color is not None:
            return f"THEME:{getattr(theme_color, 'name', str(theme_color))}"
    except Exception:
        pass
    return None

def get_fill_info(shape) -> Dict[str, Any]:
    info: Dict[str, Any] = {"type": None, "fore_color_rgb": None, "back_color_rgb": None, "transparency": None}
    try:
        fill = getattr(shape, "fill", None)
        fore = getattr(fill, "fore_color", None) if fill is not None else None
        back = getattr(fill, "back_color", None) if fill is not None else None
        info["fore_color_rgb"] = rgb_color_to_hex(fore)
        info["back_color_rgb"] = rgb_color_to_hex(back)
        if fore is not None:
            info["transparency"] = getattr(fore, "transparency", None)
        info["type"] = "solid-or-theme" if (info["fore_color_rgb"] or info["back_color_rgb"]) else "none-or-unsupported"
    except Exception:
        pass
    return info

def get_line_info(shape) -> Dict[str, Any]:
    info: Dict[str, Any] = {"width_pt": None, "color_rgb": None, "dash_style": None}
    try:
        line = getattr(shape, "line", None)
        w = getattr(line, "width", None) if line is not None else None
        if w is not None:
            try:
                info["width_pt"] = emu_to_pt(w)
            except Exception:
                info["width_pt"] = None
        color = getattr(line, "color", None) if line is not None else None
        info["color_rgb"] = rgb_color_to_hex(color)
        dash = getattr(line, "dash_style", None) if line is not None else None
        if dash is not None:
            info["dash_style"] = getattr(dash, "name", str(dash))
    except Exception:
        pass
    return info

def get_placeholder_info(shape) -> Optional[Dict[str, Any]]:
    try:
        if getattr(shape, "is_placeholder", False):
            phf = getattr(shape, "placeholder_format", None)
            return {"idx": getattr(phf, "idx", None), "type": getattr(getattr(phf, "type", None), "name", None)}
    except Exception:
        pass
    return None

def get_hyperlinks_from_text_frame(text_frame) -> List[Dict[str, Any]]:
    links: List[Dict[str, Any]] = []
    try:
        for p in getattr(text_frame, "paragraphs", []):
            for r in getattr(p, "runs", []):
                try:
                    hl = getattr(r, "hyperlink", None)
                    if hl is not None:
                        addr = getattr(hl, "address", None)
                        subaddr = getattr(hl, "sub_address", None)
                        if addr or subaddr:
                            links.append({"text": getattr(r, "text", ""), "address": addr, "sub_address": subaddr})
                except Exception:
                    continue
    except Exception:
        pass
    return links

def extract_common_shape_props(shape, order_index: int) -> Dict[str, Any]:
    position = {
        "left_emu": getattr(shape, "left", None),
        "top_emu": getattr(shape, "top", None),
        "width_emu": getattr(shape, "width", None),
        "height_emu": getattr(shape, "height", None),
    }
    pos_in = {
        "left_in": emu_to_in(position["left_emu"]),
        "top_in": emu_to_in(position["top_emu"]),
        "width_in": emu_to_in(position["width_emu"]),
        "height_in": emu_to_in(position["height_emu"]),
    }
    data: Dict[str, Any] = {
        "id": getattr(shape, "shape_id", None),
        "name": getattr(shape, "name", None),
        "type": getattr(getattr(shape, "shape_type", None), "name", None),
        "order": order_index,
        "position": {**position, **pos_in},
        "rotation_deg": getattr(shape, "rotation", None),
        "fill": get_fill_info(shape),
        "line": get_line_info(shape),
        "placeholder": get_placeholder_info(shape),
    }
    try:
        click_action = getattr(shape, "click_action", None)
        if click_action is not None:
            hl = getattr(click_action, "hyperlink", None)
            if hl is not None:
                data["hyperlink"] = {"address": getattr(hl, "address", None), "sub_address": getattr(hl, "sub_address", None)}
    except Exception:
        pass
    return data

def extract_text_shape(shape, order_index: int) -> Dict[str, Any]:
    data = extract_common_shape_props(shape, order_index)
    data["text"] = {}
    try:
        tf = getattr(shape, "text_frame", None)
        if tf is None:
            return data
        tf_info = {
            "auto_size": getattr(tf, "auto_size", None),
            "word_wrap": getattr(tf, "word_wrap", None),
            "margins_pt": {
                "left": emu_to_pt(getattr(tf, "margin_left", None)),
                "right": emu_to_pt(getattr(tf, "margin_right", None)),
                "top": emu_to_pt(getattr(tf, "margin_top", None)),
                "bottom": emu_to_pt(getattr(tf, "margin_bottom", None)),
            },
        }
        paragraphs = []
        for p in getattr(tf, "paragraphs", []):
            p_info: Dict[str, Any] = {
                "alignment": getattr(getattr(p, "alignment", None), "name", None),
                "level": getattr(p, "level", None),
                "line_spacing": getattr(p, "line_spacing", None),
                "runs": [],
            }
            for r in getattr(p, "runs", []):
                font = getattr(r, "font", None)
                p_info["runs"].append({
                    "text": getattr(r, "text", ""),
                    "font": {
                        "name": getattr(font, "name", None) if font is not None else None,
                        "size_pt": getattr(getattr(font, "size", None), "pt", None) if font is not None else None,
                        "bold": getattr(font, "bold", None) if font is not None else None,
                        "italic": getattr(font, "italic", None) if font is not None else None,
                        "underline": getattr(font, "underline", None) if font is not None else None,
                        "color_rgb": rgb_color_to_hex(getattr(font, "color", None) if font is not None else None),
                        "theme_color": getattr(getattr(font, "color", None), "theme_color", None) if font is not None else None,
                    },
                })
            paragraphs.append(p_info)
        data["text"] = {"text_frame": tf_info, "paragraphs": paragraphs, "hyperlinks": get_hyperlinks_from_text_frame(tf)}
    except Exception:
        pass
    return data

def extract_table_shape(shape, order_index: int) -> Dict[str, Any]:
    data = extract_common_shape_props(shape, order_index)
    try:
        table = getattr(shape, "table", None)
        if table is None:
            return data
        rows = len(table.rows)
        cols = len(table.columns)
        cells: List[Dict[str, Any]] = []
        for r in range(rows):
            for c in range(cols):
                try:
                    cell = table.cell(r, c)
                except Exception:
                    continue
                cell_info: Dict[str, Any] = {"row": r, "col": c, "text": getattr(cell, 'text', None)}
                try:
                    cell_info['fill_color_rgb'] = rgb_color_to_hex(getattr(getattr(cell, 'fill', None), 'fore_color', None))
                except Exception:
                    cell_info['fill_color_rgb'] = None
                try:
                    tf = getattr(cell, 'text_frame', None)
                    para_infos: List[Dict[str, Any]] = []
                    if tf is not None:
                        for p in getattr(tf, 'paragraphs', []):
                            p_info = {"alignment": getattr(getattr(p, 'alignment', None), 'name', None), "level": getattr(p, 'level', None), "runs": []}
                            for run in getattr(p, 'runs', []):
                                f = getattr(run, 'font', None)
                                p_info['runs'].append({
                                    'text': getattr(run, 'text', ''),
                                    'font': {
                                        'name': getattr(f, 'name', None) if f else None,
                                        'size_pt': getattr(getattr(f, 'size', None), 'pt', None) if f else None,
                                        'bold': getattr(f, 'bold', None) if f else None,
                                        'italic': getattr(f, 'italic', None) if f else None,
                                        'underline': getattr(f, 'underline', None) if f else None,
                                        'color_rgb': rgb_color_to_hex(getattr(f, 'color', None) if f else None),
                                    },
                                })
                            para_infos.append(p_info)
                    cell_info['paragraphs'] = para_infos
                except Exception:
                    pass
                cells.append(cell_info)
        data['table'] = {'rows': rows, 'cols': cols, 'cells': cells}
    except Exception:
        pass
    return data

def extract_picture_shape(shape, order_index: int) -> Dict[str, Any]:
    data = extract_common_shape_props(shape, order_index)
    try:
        img = getattr(shape, 'image', None)
        image_meta = {}
        if img is not None:
            image_meta['filename'] = getattr(img, 'filename', None)
            image_meta['content_type'] = getattr(img, 'content_type', None)
        data['picture'] = {
            'display_size': {
                'width_emu': getattr(shape, 'width', None),
                'height_emu': getattr(shape, 'height', None),
                'width_in': emu_to_in(getattr(shape, 'width', None)),
                'height_in': emu_to_in(getattr(shape, 'height', None)),
            },
            'image_meta': image_meta
        }
    except Exception:
        pass
    return data

def extract_chart_shape(shape, order_index: int) -> Dict[str, Any]:
    data = extract_common_shape_props(shape, order_index)
    try:
        chart = getattr(shape, 'chart', None)
        if chart is None:
            return data
        chart_info: Dict[str, Any] = {}
        try:
            chart_info['chart_type'] = getattr(getattr(chart, 'chart_type', None), 'name', None)
        except Exception:
            chart_info['chart_type'] = None
        series_list: List[Dict[str, Any]] = []
        try:
            for s in getattr(chart, 'series', []):
                s_info: Dict[str, Any] = {'name': getattr(s, 'name', None)}
                try:
                    s_info['values'] = [v.value for v in getattr(s, 'values', [])]
                except Exception:
                    s_info['values'] = None
                try:
                    s_info['categories'] = [getattr(cat, 'label', None) for cat in getattr(s, 'categories', [])]
                except Exception:
                    s_info['categories'] = None
                series_list.append(s_info)
        except Exception:
            pass
        chart_info['series'] = series_list
        data['chart'] = chart_info
    except Exception:
        pass
    return data

def extract_group_shape(shape, order_index: int) -> Dict[str, Any]:
    data = extract_common_shape_props(shape, order_index)
    children: List[Dict[str, Any]] = []
    try:
        for idx, child in enumerate(getattr(shape, 'shapes', [])):
            children.append(extract_shape(child, idx))
    except Exception:
        pass
    data['group'] = {'children': children}
    return data

def extract_generic_shape(shape, order_index: int) -> Dict[str, Any]:
    return extract_common_shape_props(shape, order_index)

def extract_shape(shape, order_index: int) -> Dict[str, Any]:
    try:
        if getattr(shape, 'has_text_frame', False):
            return extract_text_shape(shape, order_index)
    except Exception:
        pass
    try:
        st = getattr(shape, 'shape_type', None)
        if st == MSO_SHAPE_TYPE.PICTURE:
            return extract_picture_shape(shape, order_index)
        if st == MSO_SHAPE_TYPE.TABLE:
            return extract_table_shape(shape, order_index)
        if st == MSO_SHAPE_TYPE.CHART:
            return extract_chart_shape(shape, order_index)
        if st == MSO_SHAPE_TYPE.GROUP:
            return extract_group_shape(shape, order_index)
    except Exception:
        pass
    return extract_generic_shape(shape, order_index)

def extract_slide(slide, index: int) -> Dict[str, Any]:
    slide_info: Dict[str, Any] = {'index': index, 'name': getattr(slide, 'name', None), 'notes': None, 'shapes': []}
    try:
        for order_index, shape in enumerate(getattr(slide, 'shapes', [])):
            slide_info['shapes'].append(extract_shape(shape, order_index))
    except Exception:
        pass
    # notes best-effort
    try:
        notes_slide = getattr(slide, 'notes_slide', None)
        if notes_slide is not None:
            tf = getattr(notes_slide, 'notes_text_frame', None)
            if tf is not None:
                texts = []
                for p in getattr(tf, 'paragraphs', []):
                    try:
                        texts.append(''.join(getattr(run, 'text', '') for run in getattr(p, 'runs', [])))
                    except Exception:
                        pass
                slide_info['notes'] = '\n'.join([t for t in texts if t])
    except Exception:
        pass
    return slide_info

def extract_presentation(prs: Presentation) -> Dict[str, Any]:
    top: Dict[str, Any] = {'presentation': {}, 'slides': []}
    try:
        width = getattr(prs, 'slide_width', None)
        height = getattr(prs, 'slide_height', None)
        top['presentation']['slide_size'] = {
            'width_emu': width,
            'height_emu': height,
            'width_in': emu_to_in(width),
            'height_in': emu_to_in(height),
        }
    except Exception:
        pass
    try:
        for idx, slide in enumerate(getattr(prs, 'slides', [])):
            top['slides'].append(extract_slide(slide, idx))
    except Exception:
        pass
    return top

def iter_all_shapes(slide):
    try:
        for shape in getattr(slide, 'shapes', []):
            if getattr(shape, 'shape_type', None) == MSO_SHAPE_TYPE.GROUP:
                for child in getattr(shape, 'shapes', []):
                    for inner in _iter_group(child):
                        yield inner
            else:
                yield shape
    except Exception:
        return

def _iter_group(shape):
    try:
        if getattr(shape, 'shape_type', None) == MSO_SHAPE_TYPE.GROUP:
            for s in getattr(shape, 'shapes', []):
                for inner in _iter_group(s):
                    yield inner
        else:
            yield shape
    except Exception:
        return
# ==== End of helpers ====
"""


@dataclass
class VerifierResult:
    passed: bool
    reason: str
    llm_code_path: str
    stdout: str
    stderr: str

    def to_dict(self) -> Dict[str, Any]:
        return {
            "passed": self.passed,
            "reason": self.reason,
            "llm_code_path": self.llm_code_path,
            "stdout": self.stdout,
            "stderr": self.stderr,
        }


def _load_api_key() -> str:
    key = os.getenv("IRA_API_KEY") or os.getenv("OPENAI_API_KEY") or os.getenv("deerapi_key")
    if not key:
        try:
            from dotenv import load_dotenv
            load_dotenv()
            key = os.getenv("IRA_API_KEY") or os.getenv("OPENAI_API_KEY") or os.getenv("deerapi_key")
        except Exception:
            pass
    if not key:
        raise RuntimeError("Missing API key. Set IRA_API_KEY in the environment or a .env file.")
    return key


def _build_client() -> Any:
    api_key = _load_api_key()
    import openai  # type: ignore
    client = openai.OpenAI(api_key=api_key, base_url=API_URL)
    return client


def _call_llm_via_client(system: str, user: str) -> str:
    client = _build_client()
    resp = client.chat.completions.create(
        model=MODEL_NAME,
        messages=[
            {"role": "system", "content": system},
            {"role": "user", "content": user},
        ],
    )
    content = getattr(resp.choices[0].message, "content", None)
    if isinstance(content, str):
        return content
    try:
        return resp.choices[0].message["content"]  # type: ignore[index]
    except Exception:
        return ""


def _call_llm_via_http(system: str, user: str) -> str:
    api_key = _load_api_key()
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {api_key}",
    }
    payload = {
        "model": MODEL_NAME,
        "messages": [
            {"role": "system", "content": system},
            {"role": "user", "content": user},
        ],
        "temperature": 0,
    }
    last_err = None
    url = f"{API_URL}chat/completions"
    for i in range(1, HTTP_RETRY + 1):
        try:
            try:
                import requests  # type: ignore
                r = requests.post(url, headers=headers, json=payload, timeout=HTTP_TIMEOUT_SEC)
                r.raise_for_status()
                data = r.json()
            except Exception as e_req:
                from urllib.request import Request, urlopen
                import json as _json
                try:
                    req = Request(url, data=_json.dumps(payload).encode("utf-8"), headers=headers, method="POST")
                    with urlopen(req, timeout=HTTP_TIMEOUT_SEC) as resp:
                        txt = resp.read().decode("utf-8")
                    data = _json.loads(txt)
                except Exception as e_url:
                    last_err = e_url if e_req is None else e_req
                    raise last_err
            try:
                return data["choices"][0]["message"]["content"]
            except Exception as e:
                last_err = e
        except Exception as e:
            last_err = e
        if i < HTTP_RETRY:
            time.sleep(2 * i)
    if last_err:
        raise last_err
    return ""


def _llm_prompt(checks_text: str, error_context: Optional[str] = None) -> Tuple[str, str]:
    system = (
        "You are a careful Python auditor that writes a single, self-contained Python script "
        "to verify whether a set of simple checks are satisfied on a PPTX file. "
        "Output only one Python code block fenced by triple backticks. Do not add any extra text. "
        "Your script must:")
    system += (
        "\n- Use python-pptx (pptx module) to read the PPTX file in read-only manner."
        "\n- Implement: def check(pptx_path: str, checks: str) -> dict with keys: passed (bool), details (str). The 'checks' is a simple string of atomic items (e.g., newline or semicolon separated) that the script should evaluate."
        "\n- Provide a __main__ that parses two required CLI args: --ppt and --checks; then prints exactly one JSON line to stdout (no extra text)."
        "\n- Never modify files. No network. No os/subprocess/requests/shutil/socket/urllib usage. No eval/exec. Avoid sys.exit unless necessary."
        "\n- Choose which slides/shapes to inspect strictly from the provided checks text. Treat the checks as atomic items and evaluate them conservatively. If the checks don't specify a specific slide, you may check all slides and state your assumption in details."
        "\n- Be robust to text in text frames and text inside table cells; iterate runs and inspect formatting where relevant to the task."
        "\n- If the required formatting properties are ambiguous, provide a conservative judgement and explain it in details."
        "\n- python-pptx tips: iterate as: for slide in prs.slides: for shape in slide.shapes: if shape.has_text_frame: ...; if shape.has_table: for row in shape.table.rows: for cell in row.cells: tf = cell.text_frame; ... Do not use slide.has_text_frame (it doesn't exist)."
        "\n- Underline: Prefer boolean detection: treat underlined as (run.font.underline is True). If you strictly need enum, import from pptx.enum.text import MSO_UNDERLINE and compare under try/except only."
        "\n- To read RGB from run.font.color.rgb, use tuple indexing: r,g,b = (rgb[0], rgb[1], rgb[2]). Do not use .red/.green/.blue attributes."
        "\n- Group shapes: if shape.shape_type == MSO_SHAPE_TYPE.GROUP, iterate shape.shapes recursively to reach inner text/table shapes."
        "\n- Shapes enum import path note: from pptx.enum.shapes import MSO_SHAPE_TYPE (do NOT import from pptx.enum.shape)."
        "\n- Color note: run.font.color may be theme/auto. If color.rgb is None, explain inability to judge precisely or apply a conservative assumption and state it in details."
        "\n- Your script must never crash; wrap main execution in try/except and always print one JSON line with keys 'passed' and 'details'."
    )

    system += (
        "\n\nNote: Predefined helper functions for robust PPTX extraction are injected at the top of the script and available to call: "
        "extract_presentation, extract_slide, extract_shape, extract_text_shape, extract_table_shape, rgb_color_to_hex, iter_all_shapes, get_fill_info, get_line_info, get_placeholder_info. "
        "Use these helpers to focus on the task-specific verification logic."
    )

    system += (
        "\n\npython-pptx Quick Reference:" \
        "\n- Imports:" \
        "\n  from pptx import Presentation" \
        "\n  from pptx.enum.shapes import MSO_SHAPE_TYPE  # note: enum.shapes (not enum.shape)" \
        "\n- Group traversal:" \
        "\n  def iter_shapes(shape):" \
        "\n      from pptx.enum.shapes import MSO_SHAPE_TYPE" \
        "\n      if shape.shape_type == MSO_SHAPE_TYPE.GROUP:" \
        "\n          for s in shape.shapes: yield from iter_shapes(s)" \
        "\n      else: yield shape" \
        "\n- Text & tables:" \
        "\n  if shape.has_text_frame: tf = shape.text_frame; for p in tf.paragraphs: for run in p.runs: ..." \
        "\n  if shape.has_table: for row in shape.table.rows: for cell in row.cells: tf = cell.text_frame; ..." \
        "\n- Colors:" \
        "\n  rgb = run.font.color.rgb  # may be None for theme/auto; if None, explain limitation and judge conservatively" \
        "\n  r,g,b = rgb[0], rgb[1], rgb[2]  # do NOT use .red/.green/.blue" \
        "\n- Output & safety:" \
        "\n  Print exactly one JSON line to stdout with keys 'passed' and 'details'. No file writes, no network, no os/subprocess/requests/shutil/socket/urllib, no eval/exec."
    )

    user = (
        "Write the script now. It must be a single code block:"\
        "\n```python\n# your script here\n```\n\n"\
        "The checks to verify are:\n" + checks_text + "\n"\
        "The script will be executed as: python ppt_checker.py --ppt /abs/path.pptx --checks '<CHECKS>'\n"\
        "Remember to print only a single JSON line on stdout with keys: passed, details."
    )
    if error_context:
        user += ("\n\nThe previous attempt failed with the following error/output. "
                 "Regenerate a corrected script that avoids these issues and adheres to the rules.\n" + error_context)
    return system, user


def _extract_code_block(text: str) -> Optional[str]:
    pattern = r"```(?:python)?\n([\s\S]*?)\n```"
    m = re.search(pattern, text, re.IGNORECASE)
    if not m:
        return None
    return m.group(1)


def _basic_static_safety_check(code: str) -> Optional[str]:
    forbidden_imports = ["os", "subprocess", "shutil", "socket", "requests", "urllib"]
    for name in forbidden_imports:
        if re.search(rf"^\s*(import|from)\s+{re.escape(name)}\b", code, re.MULTILINE):
            return f"Forbidden import detected: {name}"
    if re.search(r"\beval\s*\(", code):
        return "Forbidden call: eval()"
    if re.search(r"\bexec\s*\(", code):
        return "Forbidden call: exec()"
    if re.search(r"open\s*\(.*['\"]\s*[wax]\s*['\"]", code):
        return "Forbidden file write mode detected"
    if re.search(r"http[s]?://", code):
        return "Forbidden network usage detected"
    if re.search(r"from\s+pptx\.enum\.shape\b", code):
        return "Incorrect import path: use 'from pptx.enum.shapes import MSO_SHAPE_TYPE'"
    if re.search(r"\bMSO_TEXT_UNDERLINE\b", code):
        return "Invalid enum MSO_TEXT_UNDERLINE: rely on run.font.underline or MSO_UNDERLINE from pptx.enum.text"
    return None


def _write_generated_checker(code: str, target_path: Optional[str] = None) -> str:
    path = target_path or GENERATED_CHECKER_PATH
    os.makedirs(os.path.dirname(path), exist_ok=True)
    prelude = EXTRACT_PRELUDE.strip() + "\n\n"
    with open(path, "w", encoding="utf-8") as f:
        f.write(prelude + code)
    return path


def _build_checker_cmd(checker_path: str, ppt_path: str, checks_text: str) -> str:
    return f"{shlex.quote(sys.executable)} {shlex.quote(checker_path)} --ppt {shlex.quote(ppt_path)} --checks {shlex.quote(checks_text)}"


def _run_generated_checker(checker_path: str, ppt_path: str, checks_text: str, timeout_sec: int = 20) -> Tuple[str, str, int, str]:
    cmd = _build_checker_cmd(checker_path, ppt_path, checks_text)
    proc = subprocess.run(
        cmd,
        shell=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        timeout=timeout_sec,
        text=True,
    )
    return proc.stdout.strip(), proc.stderr.strip(), proc.returncode, cmd


def _parse_json_line(s: str) -> Optional[Dict[str, Any]]:
    line = s.strip().splitlines()[0] if s.strip() else ""
    candidate = line
    if not candidate or not candidate.strip().startswith("{"):
        m = re.search(r"\{[\s\S]*\}", s)
        candidate = m.group(0) if m else ""
    if not candidate:
        return None
    try:
        return json.loads(candidate)
    except Exception:
        return None


def _ensure_dir(path: str) -> None:
    try:
        os.makedirs(path, exist_ok=True)
    except Exception:
        pass


def _safe_write(path: str, content: str) -> None:
    try:
        _ensure_dir(os.path.dirname(path))
        with open(path, "w", encoding="utf-8") as f:
            f.write(content)
    except Exception:
        pass


def _log(trace: bool, msg: str) -> None:
    if trace:
        try:
            sys.stderr.write(msg.rstrip("\n") + "\n")
            sys.stderr.flush()
        except Exception:
            pass


def run_o3_verifier_ppt(checks_text: str, ppt_path: str, attempts: int = 3, trace: bool = False, log_root_dir: Optional[str] = None, output_dir: Optional[str] = None) -> VerifierResult:
    if not checks_text:
        raise ValueError("checks_text is required")
    if not ppt_path:
        raise ValueError("ppt_path is required")
    if not os.path.isabs(ppt_path):
        ppt_path = os.path.abspath(ppt_path)
    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"PPT file not found: {ppt_path}")

    ts = time.strftime("%Y%m%d-%H%M%S")
    if output_dir:
        base = os.path.join(output_dir, "doc_tools")
        log_root_dir = log_root_dir or os.path.join(base, "attempt_logs", f"ppt-{ts}")
        checker_path = os.path.join(base, "ppt_checker.py")
    else:
        if log_root_dir is None:
            log_root_dir = os.path.join(tempfile.gettempdir(), "ira_attempt_logs", f"ppt-{ts}")
        checker_path = GENERATED_CHECKER_PATH
    _ensure_dir(log_root_dir)

    summary_path = os.path.join(log_root_dir, "summary.txt")

    attempt = 1
    error_ctx: Optional[str] = None
    last_stdout = ""
    last_stderr = ""
    last_code_path = ""
    last_content = ""

    while attempt <= max(1, attempts):
        attempt_dir = os.path.join(log_root_dir, f"attempt-{attempt}")
        _ensure_dir(attempt_dir)
        _log(trace, f"[ppt attempt {attempt}] starting...")

        system, user = _llm_prompt(checks_text, error_context=error_ctx)
        _safe_write(os.path.join(attempt_dir, "00-prompt-system.txt"), system)
        _safe_write(os.path.join(attempt_dir, "01-prompt-user.txt"), user)

        content = ""
        err_sdk: Optional[Exception] = None
        used_path = "sdk"
        try:
            content = _call_llm_via_client(system, user)
        except Exception as e:
            err_sdk = e

        if not content or not _extract_code_block(content):
            used_path = "http"
            try:
                content = _call_llm_via_http(system, user)
            except Exception as e:
                meta = [
                    f"used_path={used_path}",
                    f"sdk_error={repr(err_sdk)}",
                    f"http_error={repr(e)}",
                ]
                _safe_write(os.path.join(attempt_dir, "02b-llm-meta.txt"), "\n".join(meta))
                reason = f"LLM call failed (sdk={err_sdk}) (http={e})"
                if attempt >= max(1, attempts):
                    _safe_write(summary_path, (open(summary_path, "r", encoding="utf-8").read() + "\n") if os.path.exists(summary_path) else "")
                    _safe_write(summary_path, f"attempt {attempt}: LLM call failed. {reason}\n")
                    _log(trace, f"[ppt attempt {attempt}] LLM call failed. {reason}")
                    return VerifierResult(
                        passed=False,
                        reason=reason,
                        llm_code_path=last_code_path,
                        stdout=str(last_stdout or content),
                        stderr=str(last_stderr),
                    )
                else:
                    error_ctx = f"Previous LLM call failed. Reason: {reason}"
                    _safe_write(os.path.join(attempt_dir, "02-llm-raw.txt"), str(content))
                    _safe_write(os.path.join(attempt_dir, "02b-llm-meta.txt"), "\n".join(meta))
                    _safe_write(summary_path, (open(summary_path, "r", encoding="utf-8").read() + "\n") if os.path.exists(summary_path) else "")
                    _safe_write(summary_path, f"attempt {attempt}: LLM call failed. Will retry.\n")
                    _log(trace, f"[ppt attempt {attempt}] LLM call failed, retrying...")
                    attempt += 1
                    continue

        last_content = content
        _safe_write(os.path.join(attempt_dir, "02-llm-raw.txt"), str(content))
        _safe_write(os.path.join(attempt_dir, "02b-llm-meta.txt"), f"used_path={used_path}\nsdk_error={repr(err_sdk)}")

        code = _extract_code_block(content)
        if not code:
            _safe_write(os.path.join(attempt_dir, "03-extracted_code.py"), "")
            if attempt >= max(1, attempts):
                _safe_write(summary_path, (open(summary_path, "r", encoding="utf-8").read() + "\n") if os.path.exists(summary_path) else "")
                _safe_write(summary_path, f"attempt {attempt}: no valid code block.\n")
                _log(trace, f"[ppt attempt {attempt}] no valid code block; stop.")
                return VerifierResult(
                    passed=False,
                    reason="Failed to extract code block from LLM response",
                    llm_code_path=last_code_path,
                    stdout=content,
                    stderr="",
                )
            else:
                error_ctx = "The earlier response did not include a valid ```python code block. Ensure you wrap the entire script in triple backticks."
                _safe_write(summary_path, (open(summary_path, "r", encoding="utf-8").read() + "\n") if os.path.exists(summary_path) else "")
                _safe_write(summary_path, f"attempt {attempt}: no code block; retry.\n")
                _log(trace, f"[ppt attempt {attempt}] no code block; retrying...")
                attempt += 1
                continue

        _safe_write(os.path.join(attempt_dir, "03-extracted_code.py"), code)

        unsafe_reason = _basic_static_safety_check(code)
        if unsafe_reason:
            _safe_write(os.path.join(attempt_dir, "04-safety_check.txt"), f"REJECTED: {unsafe_reason}")
            if attempt >= max(1, attempts):
                _safe_write(summary_path, (open(summary_path, "r", encoding="utf-8").read() + "\n") if os.path.exists(summary_path) else "")
                _safe_write(summary_path, f"attempt {attempt}: safety check failed: {unsafe_reason}.\n")
                _log(trace, f"[ppt attempt {attempt}] safety check failed; stop.")
                return VerifierResult(
                    passed=False,
                    reason=f"Generated code rejected by safety check: {unsafe_reason}",
                    llm_code_path=last_code_path,
                    stdout=code,
                    stderr="",
                )
            else:
                error_ctx = f"Your previous code failed safety checks: {unsafe_reason}. Remove the offending imports/calls and regenerate."
                _safe_write(summary_path, (open(summary_path, "r", encoding="utf-8").read() + "\n") if os.path.exists(summary_path) else "")
                _safe_write(summary_path, f"attempt {attempt}: safety check failed; retry.\n")
                _log(trace, f"[ppt attempt {attempt}] safety check failed; retrying...")
                attempt += 1
                continue
        else:
            _safe_write(os.path.join(attempt_dir, "04-safety_check.txt"), "OK")

        path = _write_generated_checker(code, target_path=checker_path)
        last_code_path = path
        per_attempt_copy = os.path.join(attempt_dir, f"generated_checker_attempt_{attempt}.py")
        _safe_write(per_attempt_copy, code)
        _safe_write(os.path.join(attempt_dir, "05-generated_checker_path.txt"), f"main={path}\ncopy={per_attempt_copy}")

        try:
            stdout, stderr, rc, cmd = _run_generated_checker(checker_path, ppt_path, checks_text)
            _safe_write(os.path.join(attempt_dir, "06-run_command.txt"), cmd)
            _safe_write(os.path.join(attempt_dir, "07-stdout.txt"), stdout)
            _safe_write(os.path.join(attempt_dir, "08-stderr.txt"), stderr)
        except subprocess.TimeoutExpired:
            _safe_write(os.path.join(attempt_dir, "06-run_command.txt"), _build_checker_cmd(checker_path, ppt_path, checks_text))
            _safe_write(os.path.join(attempt_dir, "08-stderr.txt"), "TimeoutExpired")
            if attempt >= max(1, attempts):
                _safe_write(summary_path, (open(summary_path, "r", encoding="utf-8").read() + "\n") if os.path.exists(summary_path) else "")
                _safe_write(summary_path, f"attempt {attempt}: checker timed out.\n")
                _log(trace, f"[ppt attempt {attempt}] checker timed out; stop.")
                return VerifierResult(
                    passed=False,
                    reason="Generated checker timed out",
                    llm_code_path=path,
                    stdout=last_stdout,
                    stderr=last_stderr,
                )
            else:
                error_ctx = "Your previous script timed out. Make the logic efficient and avoid long loops."
                _safe_write(summary_path, (open(summary_path, "r", encoding="utf-8").read() + "\n") if os.path.exists(summary_path) else "")
                _safe_write(summary_path, f"attempt {attempt}: checker timed out; retry.\n")
                _log(trace, f"[ppt attempt {attempt}] checker timed out; retrying...")
                attempt += 1
                continue
        except Exception as e:
            _safe_write(os.path.join(attempt_dir, "06-run_command.txt"), _build_checker_cmd(checker_path, ppt_path, checks_text))
            _safe_write(os.path.join(attempt_dir, "08-stderr.txt"), f"Exception: {repr(e)}")
            if attempt >= max(1, attempts):
                _safe_write(summary_path, (open(summary_path, "r", encoding="utf-8").read() + "\n") if os.path.exists(summary_path) else "")
                _safe_write(summary_path, f"attempt {attempt}: failed to run checker: {repr(e)}.\n")
                _log(trace, f"[ppt attempt {attempt}] failed to run checker; stop.")
                return VerifierResult(
                    passed=False,
                    reason=f"Failed to run generated checker: {e}",
                    llm_code_path=path,
                    stdout=last_stdout,
                    stderr=last_stderr,
                )
            else:
                error_ctx = f"Your previous script raised an exception on execution: {e}. Fix and regenerate."
                _safe_write(summary_path, (open(summary_path, "r", encoding="utf-8").read() + "\n") if os.path.exists(summary_path) else "")
                _safe_write(summary_path, f"attempt {attempt}: checker exception; retry.\n")
                _log(trace, f"[ppt attempt {attempt}] checker exception; retrying...")
                attempt += 1
                continue

        last_stdout, last_stderr = stdout, stderr
        parsed = _parse_json_line(stdout)
        if parsed:
            _safe_write(os.path.join(attempt_dir, "09-parsed_json.json"), json.dumps(parsed, ensure_ascii=False, indent=2))
        else:
            _safe_write(os.path.join(attempt_dir, "09-parsed_json.json"), "<invalid or missing JSON line>")

        if not parsed or not isinstance(parsed, dict) or "passed" not in parsed:
            short_err = (stderr or "").strip()
            if len(short_err) > 800:
                short_err = short_err[-800:]
            if attempt >= max(1, attempts):
                _safe_write(summary_path, (open(summary_path, "r", encoding="utf-8").read() + "\n") if os.path.exists(summary_path) else "")
                _safe_write(summary_path, f"attempt {attempt}: invalid JSON output.\n")
                _log(trace, f"[ppt attempt {attempt}] invalid JSON output; stop.")
                return VerifierResult(
                    passed=False,
                    reason="Generated checker did not produce a valid single-line JSON with 'passed'",
                    llm_code_path=path,
                    stdout=stdout,
                    stderr=stderr,
                )
            else:
                error_ctx = (
                    "Your script must print exactly one JSON line on stdout with keys 'passed' and 'details'. "
                    "The previous run failed. Provide a corrected script."
                )
                _safe_write(summary_path, (open(summary_path, "r", encoding="utf-8").read() + "\n") if os.path.exists(summary_path) else "")
                _safe_write(summary_path, f"attempt {attempt}: invalid JSON output; retry.\n")
                _log(trace, f"[ppt attempt {attempt}] invalid JSON; retrying...")
                attempt += 1
                continue

        passed = bool(parsed.get("passed"))
        details = str(parsed.get("details", ""))

        if not passed:
            short_details = details.strip()
            if len(short_details) > 800:
                short_details = short_details[:800]
            if attempt >= max(1, attempts):
                _safe_write(summary_path, (open(summary_path, "r", encoding="utf-8").read() + "\n") if os.path.exists(summary_path) else "")
                _safe_write(summary_path, f"attempt {attempt}: checker returned failed.\n")
                _log(trace, f"[ppt attempt {attempt}] checker returned failed; stop.")
                return VerifierResult(
                    passed=False,
                    reason=details,
                    llm_code_path=path,
                    stdout=stdout,
                    stderr=stderr,
                )
            else:
                error_ctx = (
                    "Your previous checker returned passed=false with details: " + short_details + "\n"
                    "Regenerate a corrected checker that fixes this issue while strictly following all rules."
                )
                _safe_write(summary_path, (open(summary_path, "r", encoding="utf-8").read() + "\n") if os.path.exists(summary_path) else "")
                _safe_write(summary_path, f"attempt {attempt}: checker returned failed; retry.\n")
                _log(trace, f"[ppt attempt {attempt}] checker returned failed; retrying...")
                last_stdout, last_stderr = stdout, stderr
                attempt += 1
                continue

        _safe_write(summary_path, (open(summary_path, "r", encoding="utf-8").read() + "\n") if os.path.exists(summary_path) else "")
        _safe_write(summary_path, f"attempt {attempt}: SUCCESS.\n")
        _log(trace, f"[ppt attempt {attempt}] SUCCESS. Artifacts at: {attempt_dir}")
        return VerifierResult(
            passed=True,
            reason=details,
            llm_code_path=path,
            stdout=stdout,
            stderr=stderr,
        )

    _safe_write(summary_path, (open(summary_path, "r", encoding="utf-8").read() + "\n") if os.path.exists(summary_path) else "")
    _safe_write(summary_path, "No attempt yielded a valid result.\n")
    _log(trace, f"All attempts exhausted. Logs at: {log_root_dir}")
    return VerifierResult(
        passed=False,
        reason="Exhausted attempts without a valid result",
        llm_code_path=last_code_path,
        stdout=last_stdout or last_content,
        stderr=last_stderr,
    )


class CheckPptFileTool(Tool):
    name = "checkpptfile"
    description = (
    """Validator for .pptx files using an LLM-generated python-pptx checker.
Returns pass/fail with concise, auditable reasoning.

WHAT THIS TOOL CAN DO:
- Verify explicitly specified targets (e.g., slide indices, named shapes).
- Check notes text, table cell text and font properties.
- Check chart series values and categories.
- Resolve hyperlinks via r:id relationships.
- Traverse grouped shapes.
- Inspect placeholder info.
- Read image metadata and display sizes (NOT image content).
- Verify slide size, EMU → inch/pt conversions.
- Check border, fill, line styles.
- Check font and paragraph properties (bold, italic, underline, alignment, size).
- Aggregate results across specified slides and return a single pass/fail decision.

WHAT THIS TOOL CANNOT DO:
- Discover which slides, shapes, or images satisfy a condition.
- Identify visual or semantic concepts (e.g., people, faces, objects, scenes).
- Decide targets from vague checks such as “any slide with people”.

CRITICAL USAGE RULE:
- This tool is used ONLY for verification, NEVER for discovery.
- All targets MUST be explicitly specified in the checks.
- Checks relying on implicit discovery MUST be rejected or marked as NOT verifiable.

Visual or semantic discovery (e.g., locating slides with people) MUST be done
beforehand using visual observation tools (e.g., observe_current_state).
Only the resulting explicit targets may be verified here.

Limitations:
- Constrained by python-pptx; complex animations, transitions, and some low-level
  XML details are evaluated conservatively.

Inputs:
- file_path: host path to the .pptx file
- things_to_check: newline/semicolon-separated atomic checks with explicit targets

Output:
- JSON: passed, reason, llm_code_path, stdout, stderr
- Checker code and logs are recorded for audit

Prefer this tool for cross-slide or cross-part verification of EXPLICIT targets.
Do NOT use it for visual discovery or content recognition."""
)

    
    inputs = {
        "file_path": {
            "description": (
                "Path to the .pptx file on the Host (absolute or relative). "
                "Do not pass VM paths. If the file is in the VM, call get_vm_file first and then pass the returned Host path."
            ),
            "type": "string",
        },
        "things_to_check": {
            "description": "A simple string describing atomic checks to verify (e.g., newline/semicolon separated).",
            "type": "string",
        },
    }
    output_type = "string"

    def __init__(self, output_dir: Optional[str] = None):
        super().__init__()
        self.output_dir = output_dir

    def set_output_dir(self, output_dir: str) -> None:
        self.output_dir = output_dir

    def forward(self, file_path: str, things_to_check: str) -> str:
        return self.__call__(file_path, things_to_check)

    def __call__(self, file_path: str, things_to_check: str) -> str:
        try:
            if not file_path:
                return json.dumps({"passed": False, "reason": "file_path is required"}, ensure_ascii=False)
            if not things_to_check:
                return json.dumps({"passed": False, "reason": "things_to_check is required"}, ensure_ascii=False)
            abs_path = os.path.abspath(file_path)
            if not os.path.exists(abs_path):
                return json.dumps({
                    "passed": False,
                    "reason": (
                        f"File not found on host: {abs_path}. "
                        "This tool requires a Host file path. If your file is inside the VM (e.g., /home/user/...), "
                        "please use get_vm_file to download it to the Host and pass the returned Host path."
                    )
                }, ensure_ascii=False)

            res = run_o3_verifier_ppt(things_to_check, abs_path, attempts=3, trace=False, log_root_dir=None, output_dir=self.output_dir)
            return json.dumps(res.to_dict(), ensure_ascii=False)
        except Exception as e:
            return json.dumps({"passed": False, "reason": f"Error: {e}"}, ensure_ascii=False)

    def to_code_prompt(self) -> str:
        return (
            "def checkpptfile(file_path: str, things_to_check: str) -> str:\n"
            "    '''Verify a PowerPoint (.pptx) file according to a simple checks string.\n"
            "    Note: file_path must be a Host path; if the file is in the VM, first use get_vm_file to download it.\n"
            "    Returns a JSON string with keys passed, reason, llm_code_path, stdout, stderr.'''\n"
        )
